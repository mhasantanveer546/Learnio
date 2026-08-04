from pptx import Presentation


def extract_text_from_pptx(filepath):
    prs = Presentation(filepath)
    parts = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs)
                    if line.strip():
                        slide_text.append(line)

        if slide_text:
            # Slide number prefix matters here — Phase 4/5 will want to
            # cite "Slide 12" when answering questions, not just dump
            # all text with no positional context.
            parts.append(f"[Slide {i}]\n" + "\n".join(slide_text))

    extracted = "\n\n".join(parts).strip()

    if not extracted:
        raise ValueError("No extractable text found in this presentation.")

    return extracted